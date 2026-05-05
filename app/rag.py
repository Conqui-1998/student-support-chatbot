import os
import faiss
import numpy as np
from openai import OpenAI
from dotenv import load_dotenv
import re
from app.moodle_sync import has_moodle_access, is_module_enabled

try:
    from pypdf import PdfReader
except Exception:  # pragma: no cover - optional dependency
    PdfReader = None

try:
    from docx import Document
except Exception:  # pragma: no cover - optional dependency
    Document = None

try:
    from pptx import Presentation
except Exception:  # pragma: no cover - optional dependency
    Presentation = None

load_dotenv()

client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

Data_Dir = "data"
Module_Data_Dir = os.path.join(Data_Dir, "modules")
Embed_Model = "text-embedding-3-small"
Chunk_Size = 700
Top_K = 3

documents = []
index = None
module_documents = {}
module_indexes = {}

SUPPORTED_EXTENSIONS = {".md", ".pdf", ".docx", ".pptx"}

def parse_markdown_file(path: str):
    with open(path, "r", encoding="utf-8") as f:
        raw_text = f.read()
    
    lines = raw_text.splitlines()
    
    title = None
    url = None
    category = None
    
    content_start_idx = 0
    
    for i, line in enumerate(lines):
        stripped = line.strip()
        
        if stripped.startswith("Title:"):
            title = stripped.replace("Title:", "", 1).strip()
        elif stripped.startswith("URL:"):
            url = stripped.replace("URL:", "", 1).strip()
        elif stripped.startswith("Category:"):
            category = stripped.replace("Category:", "", 1).strip()
        elif stripped == "":
            content_start_idx = i + 1
            break
            
    content = "\n".join(lines[content_start_idx:]).strip()
    
    if not title:
        title = os.path.basename(path)
        
    return {
        "title": title,
        "url": url,
        "category": category,
        "content": content
    }


def parse_pdf_file(path: str):
    if PdfReader is None:
        raise RuntimeError("pypdf is not installed")

    reader = PdfReader(path)
    pages = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            pages.append(text.strip())

    return {
        "title": os.path.basename(path),
        "url": None,
        "category": "pdf",
        "content": "\n\n".join(pages).strip(),
    }


def parse_docx_file(path: str):
    if Document is None:
        raise RuntimeError("python-docx is not installed")

    doc = Document(path)
    paragraphs = []
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            paragraphs.append(text)

    return {
        "title": os.path.basename(path),
        "url": None,
        "category": "docx",
        "content": "\n\n".join(paragraphs).strip(),
    }


def parse_pptx_file(path: str):
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed")

    pres = Presentation(path)
    slides = []
    for slide_num, slide in enumerate(pres.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                    parts.append(text)
        slide_text = "\n".join(parts).strip()
        if slide_text:
            slides.append(f"Slide {slide_num}\n{slide_text}")

    return {
        "title": os.path.basename(path),
        "url": None,
        "category": "pptx",
        "content": "\n\n".join(slides).strip(),
    }


def parse_document_file(path: str):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".md":
        return parse_markdown_file(path)
    if ext == ".pdf":
        return parse_pdf_file(path)
    if ext == ".docx":
        return parse_docx_file(path)
    if ext == ".pptx":
        return parse_pptx_file(path)
    return None

def load_documents():
    docs = []
    if not os.path.isdir(Data_Dir):
        return docs

    for filename in os.listdir(Data_Dir):
        path = os.path.join(Data_Dir, filename)
        if os.path.isfile(path) and os.path.splitext(filename)[1].lower() in SUPPORTED_EXTENSIONS:
            parsed = parse_document_file(path)
            if parsed and parsed.get("content"):
                docs.extend(chunk_text(parsed))
    return docs

def sanitize_module_key(module_key: str):
    if not module_key:
        return None
    safe = re.sub(r"[^a-zA-Z0-9_-]+", "-", module_key.strip()).strip("-")
    return safe.lower() or None

def load_module_documents(module_key: str):
    module_key = sanitize_module_key(module_key)
    if not module_key:
        return []

    module_dir = os.path.join(Module_Data_Dir, module_key)
    if not os.path.isdir(module_dir):
        return []

    docs = []
    for filename in os.listdir(module_dir):
        path = os.path.join(module_dir, filename)
        if os.path.isfile(path) and os.path.splitext(filename)[1].lower() in SUPPORTED_EXTENSIONS:
            parsed = parse_document_file(path)
            if parsed and parsed.get("content"):
                docs.extend(chunk_text(parsed))
    return docs
    
def chunk_text(doc: dict, chunk_size: int = Chunk_Size):
    chunks = []
    current = ""
    
    paragraphs = doc["content"].split("\n\n")
    
    for paragraph in paragraphs:
        paragraph = paragraph.strip()
        if not paragraph:
            continue
        if len(current) + len(paragraph) < chunk_size:
            current += paragraph + "\n\n"
        else:
            if current.strip():
                chunks.append({"title": doc["title"],"url": doc["url"], "category": doc["category"], "content": current.strip()})
            current = paragraph + "\n\n"
            
    if current.strip():
        chunks.append({"title": doc["title"], "url": doc["url"], "category": doc["category"], "content": current.strip()})
    return chunks

def embed_texts(texts):
    response = client.embeddings.create(
        model=Embed_Model,
        input=texts
    )
    return [d.embedding for d in response.data]
    
def build_index():
    global documents, index
    documents = load_documents()
    if documents:
        embeddings = embed_texts([doc["content"] for doc in documents])
        matrix = np.array(embeddings).astype("float32")
        index = faiss.IndexFlatL2(matrix.shape[1])
        index.add(matrix)
    else:
        index = None

def build_module_index(module_key: str):
    module_key = sanitize_module_key(module_key)
    if not module_key:
        return None, []

    docs = load_module_documents(module_key)
    if not docs:
        return None, []

    embeddings = embed_texts([doc["content"] for doc in docs])
    matrix = np.array(embeddings).astype("float32")
    module_index = faiss.IndexFlatL2(matrix.shape[1])
    module_index.add(matrix)
    return module_index, docs

def search(query: str, top_k: int = Top_K, module_key: str = None):
    global index, documents, module_documents, module_indexes

    module_key = sanitize_module_key(module_key)

    if module_key and is_module_enabled(module_key):
        if module_key not in module_indexes:
            module_index, docs = build_module_index(module_key)
            if module_index is not None and docs:
                module_indexes[module_key] = module_index
                module_documents[module_key] = docs
        module_index = module_indexes.get(module_key)
        docs = module_documents.get(module_key, [])
        if module_index is not None and docs:
            query_embedding = embed_texts([query])[0]
            q = np.array([query_embedding]).astype("float32")
            distances, indices = module_index.search(q, top_k)

            results = []
            for dist, idx in zip(distances[0], indices[0]):
                if idx == -1:
                    continue
                results.append({
                    "title": docs[idx]["title"],
                    "url": docs[idx]["url"],
                    "category": docs[idx]["category"],
                    "content": docs[idx]["content"],
                    "score": float(dist)
                })
            if results:
                return results

    if index is None:
        build_index()

    if index is None or not documents:
        return []

    query_embedding = embed_texts([query])[0]
    q = np.array([query_embedding]).astype("float32")
    distances, indices = index.search(q, top_k)
    
    results = []
    for dist, idx in zip(distances[0], indices[0]):
        if idx == -1:
            continue
        results.append({
            "title": documents[idx]["title"],
            "url": documents[idx]["url"],
            "category": documents[idx]["category"],
            "content": documents[idx]["content"],
            "score": float(dist)
        })
    return results
